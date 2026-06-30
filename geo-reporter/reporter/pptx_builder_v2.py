"""
PPTX Builder V2 —— 融合论证式新版演示稿(Phase D:真实内容),与旧版 pptx_builder.py 并存。
按论证链组织:封面→摘要→项目概况→地质背景→工作方法→数据综合分析→靶区→钻孔→价值评估→风险→结论。
复用旧版 PptxBuilder 样式 helper(_set_slide_bg/_add_textbox/_add_bullet_list/_add_table_shape/_add_accent_line)。
"""
from datetime import datetime
from typing import Dict, Optional

from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN

from .pptx_builder import PptxBuilder
from .categories import SearchResult, get_all_categories
from .geocoder import LocationContext


def _clip(s, n=46):
    s = str(s or "").strip().replace("\n", " ")
    return (s[:n] + "…") if len(s) > n else s


class PptxBuilderV2(PptxBuilder):
    """融合论证式新版 PPT 生成器(Phase D:真实内容)。"""

    def build_pptx_v2(self, location: LocationContext, search_results: Dict[str, SearchResult],
                      output_name: Optional[str] = None, mineral_type: str = "",
                      target_figure=None, confidence: dict = None, econ_params: dict = None) -> str:
        prs = Presentation()
        prs.slide_width = self.SLIDE_WIDTH
        prs.slide_height = self.SLIDE_HEIGHT
        sr = search_results or {}
        conf = confidence or {}
        targets = getattr(target_figure, "targets", None) or []

        area = (getattr(location, "area_name", "") or "").strip()
        if not area or area.lower() == "aoi":
            area = (getattr(location, "location_str", "") or "").strip() or "研究区"

        def sec(title):
            s = prs.slides.add_slide(prs.slide_layouts[6])
            self._set_slide_bg(s)
            self._add_textbox(s, Inches(0.8), Inches(0.45), Inches(11.7), Inches(0.9), title,
                              font_size=26, font_name=self.FONT_TITLE, color=self.COLOR_ACCENT, bold=True)
            self._add_accent_line(s, Inches(0.8), Inches(1.25), Inches(3))
            return s

        def table(s, headers, rows, top=1.7, height=None):
            if not rows:
                return
            h = height or min(5.2, 0.5 + 0.42 * len(rows))
            self._add_table_shape(s, Inches(0.8), Inches(top), Inches(11.7), Inches(h), headers, rows)

        # ===== 封面 =====
        s = prs.slides.add_slide(prs.slide_layouts[6])
        self._set_slide_bg(s)
        self._add_accent_line(s, Inches(2), Inches(2.2), Inches(9.333))
        self._add_textbox(s, Inches(2), Inches(2.5), Inches(9.333), Inches(1.2),
                          "地质勘探综合评估报告(融合版)", font_size=40, font_name=self.FONT_TITLE,
                          color=self.COLOR_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        self._add_textbox(s, Inches(2), Inches(3.7), Inches(9.333), Inches(0.8), area,
                          font_size=28, font_name=self.FONT_TITLE, color=self.COLOR_ACCENT, alignment=PP_ALIGN.CENTER)
        info = [location.location_str] if location.location_str else []
        if mineral_type:
            info.append(f"目标矿种：{mineral_type}")
        info.append(datetime.now().strftime("%Y年%m月%d日"))
        self._add_textbox(s, Inches(2), Inches(4.8), Inches(9.333), Inches(0.6), "  |  ".join(info),
                          font_size=14, color=self.COLOR_LIGHT, alignment=PP_ALIGN.CENTER)
        self._add_accent_line(s, Inches(2), Inches(5.5), Inches(9.333))

        # ===== 摘要 =====
        s = sec("摘要")
        grade = conf.get("grade") or conf.get("overall_grade")
        pts = []
        if grade:
            pts.append(f"综合成矿置信定级：{grade} 级")
        if targets:
            depths = [t.get("target_depth_m") for t in targets if t.get("target_depth_m") not in (None, "")]
            pts.append(f"三维建模圈定 {len(targets)} 个预测靶区"
                       + (f"、目标深度约 {min(depths)}–{max(depths)} m" if depths else ""))
        # 取各类别首条关键发现(前 4)
        for cat in get_all_categories():
            if cat.id == "slow_variables":
                continue
            r = sr.get(cat.id)
            if r and not getattr(r, "error", None) and r.key_findings:
                pts.append(f"{cat.chapter_title}：{_clip(r.key_findings[0], 40)}")
            if len(pts) >= 6:
                break
        rec = conf.get("recommendation")
        if rec:
            pts.append(f"建议：{_clip(rec, 44)}")
        self._add_bullet_list(s, Inches(0.9), Inches(1.7), Inches(11.5), Inches(5.2), pts or ["(摘要待补充)"], font_size=16)

        # ===== 项目概况 =====
        s = sec("第一章 · 项目概况")
        ov = [f"位置：{location.location_str}（{location.country}）", f"坐标范围：{location.coords_str}"]
        if mineral_type:
            ov.append(f"目标矿种：{mineral_type}")
        n_done = sum(1 for r in sr.values() if r and not getattr(r, "error", None))
        ov.append(f"取得成果：综合 {n_done} 类地学资料"
                  + (f"、圈定 {len(targets)} 个三维靶区" if targets else "")
                  + (f"、置信 {grade} 级" if grade else ""))
        self._add_bullet_list(s, Inches(0.9), Inches(1.7), Inches(11.5), Inches(4.5), ov, font_size=16, icon="●")

        # ===== 区域与矿区地质背景 =====
        s = sec("第二章 · 区域与矿区地质背景")
        geo = sr.get("geology")
        if geo and not getattr(geo, "error", None):
            if geo.summary:
                self._add_textbox(s, Inches(0.9), Inches(1.7), Inches(11.5), Inches(1.6), _clip(geo.summary, 160), font_size=14, color=self.COLOR_LIGHT)
            if geo.key_findings:
                self._add_bullet_list(s, Inches(0.9), Inches(3.3), Inches(11.5), Inches(3.5), [_clip(k, 52) for k in geo.key_findings[:5]], font_size=15)
        else:
            self._add_textbox(s, Inches(0.9), Inches(1.7), Inches(11.5), Inches(1), "本区地质资料待补充。", font_size=15, color=self.COLOR_LIGHT)

        # ===== 工作方法与技术 =====
        s = sec("第三章 · 工作方法与技术")
        self._add_bullet_list(s, Inches(0.9), Inches(1.7), Inches(11.5), Inches(5.0), [
            "多模态异构数据 AI 智能找矿:融合遥感蚀变、构造、物探(磁/重)、化探、InSAR 形变多源证据",
            "三维成矿建模(geo-model3d)形成深部有利度体与靶点,构成「地-空-钻」三位一体证据链",
            "靶向超弱核磁共振 AI 探测(平台预留能力,本次未接入实测数据)",
            "证据来源分层标注(子系统本地实证/直连/检索),保障可追溯与可信度核验",
        ], font_size=15)

        # ===== 数据综合分析 =====
        s = sec("第四章 · 数据综合分析")
        rows = []
        for cat in get_all_categories():
            if cat.id in ("geology", "slow_variables"):
                continue
            r = sr.get(cat.id)
            if not r or getattr(r, "error", None):
                continue
            note = (r.key_findings[0] if r.key_findings else r.summary) or ""
            if note:
                rows.append([cat.chapter_title, _clip(note, 40)])
        if rows:
            table(s, ["数据类别", "关键发现"], rows[:8])
        else:
            self._add_textbox(s, Inches(0.9), Inches(1.7), Inches(11.5), Inches(1), "本次暂无可纳入综合分析的物化探/遥感/形变证据。", font_size=15, color=self.COLOR_LIGHT)

        # ===== 隐伏矿位置评估(靶区)=====
        s = sec("第五章 · 隐伏矿位置评估")
        if targets:
            trows = [[f"#{t.get('rank','')} {t.get('grade','')}",
                      f"{t.get('longitude',0):.4f}, {t.get('latitude',0):.4f} / {t.get('target_depth_m','')}m",
                      _clip(t.get('reason',''), 34)] for t in targets[:8]]
            table(s, ["靶区(置信)", "中心坐标 / 深度", "评分理由"], trows)
        else:
            self._add_textbox(s, Inches(0.9), Inches(1.7), Inches(11.5), Inches(1), "本研究区暂无三维成矿预测靶点。", font_size=15, color=self.COLOR_LIGHT)

        # ===== 建议钻孔方案 =====
        s = sec("第六章 · 建议钻孔方案")
        holes = []
        try:
            from .synthesis import get_drill_evidence
            holes = (get_drill_evidence(location) or {}).get("holes") or []
        except Exception:
            holes = []
        if holes:
            hrows = []
            for h in holes[:8]:
                dep = h.get("target_depth_m", "")
                try:
                    kind = "斜孔" if float(dep) > 400 else "直孔"
                except (TypeError, ValueError):
                    kind = "—"
                hrows.append([h.get("hole_id", "") or f"#{h.get('rank','')}",
                              f"{h.get('lon',0):.4f}, {h.get('lat',0):.4f} · {kind}",
                              f"{dep}m / 优先级{h.get('priority', h.get('score',''))}"])
            table(s, ["计划孔", "坐标 · 孔型", "深度 / 优先级"], hrows)
        else:
            self._add_textbox(s, Inches(0.9), Inches(1.7), Inches(11.5), Inches(1.2),
                              "本次暂无 geo-drill 布孔产物;建议靶区优选后开展钻探设计(目标深度 ≤400m 用直孔,>400m 用斜孔)。",
                              font_size=15, color=self.COLOR_LIGHT)

        # ===== 资源潜力与价值评估 =====
        s = sec("第七章 · 资源潜力与价值评估")
        from .value_assessment import compute_value_assessment
        # 预算进尺优先用钻孔总进尺(上面已取 holes),无则退化为靶点深度之和
        total_m = 0.0
        for h in (holes or []):
            try:
                total_m += float(h.get("target_depth_m") or 0)
            except (TypeError, ValueError):
                pass
        if not total_m:
            total_m = sum(float(t.get("target_depth_m") or 0) for t in targets if str(t.get("target_depth_m") or "").replace(".", "").isdigit())
        va = compute_value_assessment(econ_params or {}, targets, total_m)
        if va:
            cur = va["currency"]
            vrows = [[f"{x['name']}({x['increment_pct']:.0f}%)",
                      f"新增{x['new_koz']:,.0f} / 总{x['total_koz']:,.0f} koz",
                      (f"ROI {x['roi_pct']:,.0f}% · 发现成本{x['discovery_cost_per_oz']:,.0f}{cur}/oz" if x['roi_pct'] is not None else "—")]
                     for x in va["scenarios"]]
            self._add_textbox(s, Inches(0.9), Inches(1.55), Inches(11.5), Inches(0.6),
                              f"现有资源约 {va['existing_koz']:,.0f} koz,金属价格 {va['metal_price']:,.0f} {cur}/oz", font_size=13, color=self.COLOR_LIGHT)
            table(s, ["情景", "资源增量", "价值指标"], vrows, top=2.3)
        else:
            self._add_textbox(s, Inches(0.9), Inches(1.7), Inches(11.5), Inches(1.4),
                              "定量价值评估(koz/ROI/发现成本)需上传经济参数表(资源量/品位、金属价格、钻探单价等);"
                              "上传后本页自动给出保守/基准/乐观三情景估算。", font_size=15, color=self.COLOR_LIGHT)

        # ===== 风险评估与建议 =====
        s = sec("第八章 · 风险评估与建议")
        failed = [c.chapter_title for c in get_all_categories() if sr.get(c.id) and getattr(sr[c.id], "error", None)]
        risks = ["技术风险:三维预测为多源证据综合推断,深部矿化连续性需钻探直接验证;高置信靶区优先加密",
                 "工程与经济风险:深孔施工(涌水等)与金属价格波动需实施阶段评估"]
        if failed:
            risks.insert(1, "数据风险:以下证据本次缺失/失败——" + "、".join(failed[:5]))
        self._add_bullet_list(s, Inches(0.9), Inches(1.7), Inches(11.5), Inches(5.0), risks, font_size=15)

        # ===== 结论与建议 =====
        s = sec("第九章 · 结论与建议")
        cpts = []
        if grade:
            cpts.append(f"综合成矿置信:{grade} 级")
        if targets:
            cpts.append(f"圈定 {len(targets)} 个三维预测靶区,具备进一步勘查价值")
        if rec:
            cpts.append(f"建议:{_clip(rec, 60)}")
        cpts.append("下一步:对高置信靶区优先钻探验证,并补充缺失证据闭合证据链")
        self._add_bullet_list(s, Inches(0.9), Inches(1.7), Inches(11.5), Inches(5.0), cpts, font_size=16)

        if output_name is None:
            output_name = f"{area}_{datetime.now().strftime('%Y%m%d')}"
        out = self.output_dir / f"{output_name}_新版.pptx"
        prs.save(str(out))
        return str(out)
